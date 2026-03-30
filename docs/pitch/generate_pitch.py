"""
Jarvis VC Pitch Deck Generator — v2
Startup Foundry LaunchPad Submission
Theme: Jarvis × Iron Man (terra #D4775A, gold #E09D5C, dusk #6B7FB5, canvas #1C1A17)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── Palette ──────────────────────────────────────────────────────────────────
CANVAS   = RGBColor(0x1C, 0x1A, 0x17)
CARD     = RGBColor(0x25, 0x23, 0x20)
CARD2    = RGBColor(0x20, 0x1E, 0x1B)
TERRA    = RGBColor(0xD4, 0x77, 0x5A)
GOLD     = RGBColor(0xE0, 0x9D, 0x5C)
DUSK     = RGBColor(0x6B, 0x7F, 0xB5)
SAGE     = RGBColor(0x4A, 0x7B, 0x6B)
WHITE    = RGBColor(0xFA, 0xF8, 0xF4)
MUTED    = RGBColor(0x9C, 0x94, 0x88)
DIMMED   = RGBColor(0x6B, 0x65, 0x60)

SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=CANVAS):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def tag(slide, text, l, t, color=TERRA):
    """Small uppercase label."""
    txt(slide, text, l, t, Inches(8), Inches(0.28),
        size=8, bold=True, color=color)

def hud(slide, color=TERRA):
    """Corner bracket marks."""
    m, s, lw = Inches(0.15), Inches(0.22), Inches(0.025)
    W, H = SLIDE_W, SLIDE_H
    for (x, y, w, h) in [
        (m, m, s, lw), (m, m, lw, s),
        (W-m-s, m, s, lw), (W-m-lw, m, lw, s),
        (m, H-m-lw, s, lw), (m, H-m-s, lw, s),
        (W-m-s, H-m-lw, s, lw), (W-m-lw, H-m-s, lw, s),
    ]:
        r = slide.shapes.add_shape(1, x, y, w, h)
        r.fill.solid()
        r.fill.fore_color.rgb = color
        r.line.fill.background()

def scan(slide, color=TERRA):
    """Bottom scan line."""
    r = rect(slide, Inches(1), SLIDE_H - Inches(0.05),
             SLIDE_W - Inches(2), Inches(0.03), color)

def bullet(slide, dot_color, text, l, t, w=Inches(11.5), size=13):
    d = rect(slide, l, t + Inches(0.1), Inches(0.09), Inches(0.09), dot_color)
    txt(slide, text, l + Inches(0.2), t, w - Inches(0.2), Inches(0.38),
        size=size, color=WHITE)

def card_block(slide, l, t, w, h, accent_color=None):
    rect(slide, l, t, w, h, CARD)
    if accent_color:
        rect(slide, l, t, w, Inches(0.05), accent_color)

def pill(slide, text, l, t, color=TERRA):
    """Small pill label."""
    r = rect(slide, l, t, Inches(0.06), Inches(0.32), color)
    txt(slide, text, l + Inches(0.14), t + Inches(0.02),
        Inches(3), Inches(0.28), size=9, color=color, bold=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ════════════════════════════════════════════════════════════════════════════
s1 = add_slide(); bg(s1); hud(s1); scan(s1)

# Arc reactor
arc = rect(s1, SLIDE_W/2 - Inches(0.09), Inches(0.2), Inches(0.18), Inches(0.18), GOLD)

txt(s1, "JARVIS · AI PRODUCTIVITY ENGINE",
    Inches(0), Inches(1.5), SLIDE_W, Inches(0.35),
    size=9, bold=True, color=DUSK, align=PP_ALIGN.CENTER)

txt(s1, "Your AI that thinks,",
    Inches(0), Inches(2.05), SLIDE_W, Inches(0.85),
    size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s1, "so you don't have to.",
    Inches(0), Inches(2.85), SLIDE_W, Inches(0.85),
    size=46, bold=True, color=TERRA, align=PP_ALIGN.CENTER)

rect(s1, SLIDE_W/2 - Inches(1.6), Inches(3.85), Inches(3.2), Inches(0.025), GOLD)

txt(s1, "MURALI MADHAV  ·  JARVIS ENGINE  ·  2026",
    Inches(0), Inches(4.0), SLIDE_W, Inches(0.3),
    size=10, color=MUTED, align=PP_ALIGN.CENTER)

txt(s1, "Local-first AI that converts your brain dump into a mathematically optimal,\n"
        "guilt-free daily schedule — powered by constraint solving, LLMs & behavioral science.",
    Inches(2), Inches(4.45), Inches(9.33), Inches(0.75),
    size=13, color=MUTED, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM + INSIGHT (merged)
# ════════════════════════════════════════════════════════════════════════════
s2 = add_slide(); bg(s2); hud(s2); scan(s2)

tag(s2, "01 / PROBLEM STATEMENT + YOUR INSIGHT", Inches(0.6), Inches(0.32))
txt(s2, "Students are drowning in cognitive overhead.",
    Inches(0.6), Inches(0.65), Inches(10), Inches(0.6),
    size=30, bold=True, color=WHITE)

# 3 stat boxes
for i, (val, sub, color) in enumerate([
    ("3.2 hrs/day",    "lost planning\nbefore any real work", TERRA),
    ("5 / 5",          "interviewees said:\n\"I know what to do.\nI just can't start.\"", GOLD),
    ("78%",            "feel weekly guilt\nfrom missed tasks", DUSK),
]):
    l = Inches(0.6) + i * Inches(4.1)
    card_block(s2, l, Inches(1.45), Inches(3.9), Inches(1.3), color)
    txt(s2, val,  l + Inches(0.2), Inches(1.55), Inches(3.5), Inches(0.52),
        size=28, bold=True, color=color)
    txt(s2, sub,  l + Inches(0.2), Inches(2.05), Inches(3.5), Inches(0.65),
        size=10, color=MUTED)

# Insight section
txt(s2, "FOUNDER INSIGHT",
    Inches(0.6), Inches(2.95), Inches(4), Inches(0.25),
    size=8, bold=True, color=GOLD)
txt(s2, "I lived this. As an engineering student managing exams, side projects, and self-study — "
        "I was spending more energy figuring out what to work on than actually working. "
        "Notion didn't think for me. Calendars didn't adapt. No tool decomposed my goals into "
        "an executable day automatically.",
    Inches(0.6), Inches(3.22), Inches(7.8), Inches(0.85),
    size=12, color=MUTED)

# Why tools fail
txt(s2, "Why existing tools fail:",
    Inches(0.6), Inches(4.18), Inches(5), Inches(0.28),
    size=11, bold=True, color=WHITE)
for d, color in [
    ("Notion / Todoist are empty containers — they don't think for you", TERRA),
    ("Calendars schedule but don't decompose goals into achievable steps", GOLD),
    ("AI chatbots give suggestions, not constraint-aware executable plans", DUSK),
]:
    bullet(s2, color, d, Inches(0.6),
           Inches(4.18) + [
               "Notion / Todoist are empty containers — they don't think for you",
               "Calendars schedule but don't decompose goals into achievable steps",
               "AI chatbots give suggestions, not constraint-aware executable plans",
           ].index(d) * Inches(0.42) + Inches(0.38), size=12)

# Quote
rect(s2, Inches(0.6), Inches(5.72), Inches(0.06), Inches(0.65), TERRA)
rect(s2, Inches(0.66), Inches(5.72), Inches(11.4), Inches(0.65), CARD)
txt(s2, '"I feel guilty constantly. I have everything I need — but I can\'t turn it into a plan that sticks."',
    Inches(0.85), Inches(5.79), Inches(11), Inches(0.35),
    size=12, italic=True, color=WHITE)
txt(s2, "— Engineering student, interview #3  ·  Full transcripts attached",
    Inches(0.85), Inches(6.17), Inches(11), Inches(0.22),
    size=9, color=DIMMED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — SOLUTION (visual pipeline)
# ════════════════════════════════════════════════════════════════════════════
s3 = add_slide(); bg(s3); hud(s3); scan(s3)

tag(s3, "02 / PROPOSED SOLUTION", Inches(0.6), Inches(0.32))
txt(s3, "Brain dump  →  Perfect schedule.",
    Inches(0.6), Inches(0.65), Inches(11), Inches(0.6),
    size=30, bold=True, color=WHITE)
txt(s3, "Six steps. Fully automated. No manual planning required.",
    Inches(0.6), Inches(1.3), Inches(9), Inches(0.32),
    size=13, color=MUTED)

# Pipeline steps
steps = [
    ("BRAIN DUMP",      TERRA, "You say anything:\ntasks, goals, habits,\ndeadlines, constraints"),
    ("JARVIS AI",       GOLD,  "Jarvis AI understands\nyour intent, tasks,\nhabits and constraints"),
    ("DECOMPOSER",      DUSK,  "Breaks goals into\n≤25 min micro-tasks\nwith psychology baked in"),
    ("SCHEDULER",       TERRA, "Mathematical engine\nbuilds a provably valid,\nconflict-free plan"),
    ("DRAFT REVIEW",    GOLD,  "You accept, edit,\nreject or chat —\nJarvis adapts live"),
    ("VOICE OF JARVIS", DUSK,  "Warm, human response\nsynthesised by\nJarvis AI"),
]

step_w = Inches(1.95)
step_h = Inches(2.2)
gap    = Inches(0.22)
start_l = Inches(0.5)

for i, (name, color, desc) in enumerate(steps):
    l = start_l + i * (step_w + gap)
    card_block(s3, l, Inches(1.78), step_w, step_h, color)
    # Step number
    txt(s3, f"0{i+1}", l + Inches(0.15), Inches(1.9),
        Inches(0.6), Inches(0.38), size=18, bold=True, color=color)
    txt(s3, name, l + Inches(0.15), Inches(2.3),
        step_w - Inches(0.2), Inches(0.3),
        size=9, bold=True, color=WHITE)
    txt(s3, desc, l + Inches(0.15), Inches(2.65),
        step_w - Inches(0.2), Inches(1.2),
        size=10, color=MUTED)
    # Arrow (not after last)
    if i < len(steps) - 1:
        ax = l + step_w + Inches(0.05)
        rect(s3, ax, Inches(2.75), Inches(0.1), Inches(0.03), DIMMED)

# Bottom highlight row
txt(s3, "What makes each step intelligent:",
    Inches(0.6), Inches(4.22), Inches(5), Inches(0.28),
    size=11, bold=True, color=WHITE)

for d, color in [
    ("Jarvis AI engine is local-first, privacy-first — your data never leaves your device.", TERRA),
    ("Our scheduling engine uses constraint satisfaction mathematics — produces provably valid plans, not AI guesses.", GOLD),
    ("Draft negotiation loop — you're never locked in. Edit any task, Jarvis re-solves instantly.", DUSK),
]:
    bullet(s3, color, d, Inches(0.6),
           Inches(4.22) + [
               "Jarvis AI engine is local-first, privacy-first — your data never leaves your device.",
               "Our scheduling engine uses constraint satisfaction mathematics — produces provably valid plans, not AI guesses.",
               "Draft negotiation loop — you're never locked in. Edit any task, Jarvis re-solves instantly.",
           ].index(d) * Inches(0.42) + Inches(0.38), size=12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PSYCHOLOGY ENGINE
# ════════════════════════════════════════════════════════════════════════════
s4 = add_slide(); bg(s4); hud(s4, color=DUSK); scan(s4, color=DUSK)

tag(s4, "03 / THE PSYCHOLOGY ENGINE", Inches(0.6), Inches(0.32), color=DUSK)
txt(s4, "Every task is psychologically engineered.",
    Inches(0.6), Inches(0.65), Inches(10), Inches(0.6),
    size=30, bold=True, color=WHITE)
txt(s4, "Five proven frameworks baked into the architecture — not bolted on as features.",
    Inches(0.6), Inches(1.3), Inches(10), Inches(0.32),
    size=13, color=MUTED)

frameworks = [
    ("WOOP", TERRA,
     "Wish · Outcome · Obstacle · Plan",
     "Every task chunk has an implementation intention: a concrete if-then response to the most likely obstacle. "
     "Turns vague goals into executable micro-commitments."),
    ("CLT", GOLD,
     "Cognitive Load Theory",
     "Tasks capped at ≤25 minutes. Completion criteria require active recall, not passive review. "
     "Germane load (schema-building) maximised; extraneous load (friction) minimised."),
    ("TMT", DUSK,
     "Temporal Motivation Theory",
     "Motivation = (Expectancy × Value) / (Impulsiveness × Delay). "
     "Deadlines become integer priority weights in Jarvis's scheduling engine — nearer deadline = scheduled earlier, automatically."),
    ("SM-2", TERRA,
     "Spaced Repetition (Ebbinghaus)",
     "Habits tracked with SM-2 ease factor. Each completion updates the next review interval. "
     "Memory and routine strength decay mathematically — Jarvis resurfaces them at the right moment."),
    ("Anti-Guilt", GOLD,
     "INFEASIBLE ≠ User Failure",
     "When the solver can't fit everything, Jarvis explains why — too much scope, not enough time — "
     "and guides recalibration. No guilt. No judgment. Just math."),
]

fw_w = Inches(2.45)
fw_h = Inches(2.85)
for i, (name, color, subtitle, body) in enumerate(frameworks):
    l = Inches(0.5) + i * (fw_w + Inches(0.16))
    card_block(s4, l, Inches(1.78), fw_w, fw_h, color)
    txt(s4, name,     l+Inches(0.18), Inches(1.94), fw_w-Inches(0.25), Inches(0.38),
        size=14, bold=True, color=color)
    txt(s4, subtitle, l+Inches(0.18), Inches(2.3),  fw_w-Inches(0.25), Inches(0.28),
        size=9,  bold=True, color=WHITE)
    txt(s4, body,     l+Inches(0.18), Inches(2.62), fw_w-Inches(0.25), Inches(1.85),
        size=10, color=MUTED)

txt(s4, "Result: users don't just get a schedule — they get a plan that accounts for human psychology, not just hours.",
    Inches(0.6), Inches(4.85), Inches(11.5), Inches(0.35),
    size=13, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNICAL MOAT + USPs
# ════════════════════════════════════════════════════════════════════════════
s5 = add_slide(); bg(s5); hud(s5, color=TERRA); scan(s5)

tag(s5, "04 / TECHNICAL MOAT + USPs", Inches(0.6), Inches(0.32))
txt(s5, "Why this can't be replicated in a weekend.",
    Inches(0.6), Inches(0.65), Inches(10), Inches(0.6),
    size=30, bold=True, color=WHITE)

usps = [
    ("Local-First AI",           TERRA,
     "LLMs run on-device. No data leaves the user's machine. "
     "Privacy-first by architecture — not a policy checkbox. "
     "Cloud only when explicitly needed (real-time search)."),
    ("Deterministic Scheduler",  GOLD,
     "Constraint satisfaction mathematics builds provably valid plans. "
     "Hard constraints (sleep, meetings) are enforced, not suggested. "
     "Zero AI guesses — every plan is mathematically feasible."),
    ("PEARL Memory Engine",      DUSK,
     "Jarvis observes behaviour (skips, edits, rejections) and injects "
     "inferred patterns as constraints into the solver. "
     "The longer you use it, the more it knows you."),
    ("3-Tier Memory",            TERRA,
     "Multi-tier working / recall / archival memory with decay scoring. "
     "Memories don't just change what Jarvis says — "
     "they change the mathematical constraints in the scheduler."),
    ("Self-Evolving Intents",    GOLD,
     "Embedding similarity detects novel user patterns. "
     "Gemini auto-generates an IntentBlueprint. "
     "Jarvis learns new capabilities without code changes."),
    ("Anti-Guilt Architecture",  DUSK,
     "INFEASIBLE solver result = scope problem, never user failure. "
     "Jarvis recalibrates gracefully. "
     "Designed for ADHD brains and overwhelmed students alike."),
]

usp_w = Inches(4.1)
usp_h = Inches(1.55)
for i, (name, color, body) in enumerate(usps):
    col = i % 3
    row = i // 3
    l = Inches(0.5) + col * (usp_w + Inches(0.16))
    t = Inches(1.55) + row * (usp_h + Inches(0.14))
    card_block(s5, l, t, usp_w, usp_h)
    rect(s5, l, t, Inches(0.06), usp_h, color)
    txt(s5, name, l+Inches(0.2), t+Inches(0.12), usp_w-Inches(0.3), Inches(0.3),
        size=12, bold=True, color=color)
    txt(s5, body, l+Inches(0.2), t+Inches(0.45), usp_w-Inches(0.3), Inches(1.0),
        size=10, color=MUTED)

txt(s5, "Moat: PEARL grows with every user interaction. The dataset is proprietary. The longer users stay, the harder it is to leave.",
    Inches(0.6), Inches(5.1), Inches(11.5), Inches(0.35),
    size=12, bold=True, color=TERRA)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TARGET AUDIENCE + WHY NOW
# ════════════════════════════════════════════════════════════════════════════
s6 = add_slide(); bg(s6); hud(s6, color=GOLD); scan(s6, color=GOLD)

tag(s6, "05 / TARGET AUDIENCE + WHY NOW", Inches(0.6), Inches(0.32), color=GOLD)
txt(s6, "First 100: overloaded engineering students.",
    Inches(0.6), Inches(0.65), Inches(10), Inches(0.55),
    size=28, bold=True, color=WHITE)

# Left — audience
card_block(s6, Inches(0.5), Inches(1.38), Inches(6.0), Inches(3.6), GOLD)
txt(s6, "WHO",
    Inches(0.7), Inches(1.5), Inches(5), Inches(0.25),
    size=8, bold=True, color=GOLD)
txt(s6, "Engineering & pre-professional students juggling\nexams, projects, internships — simultaneously.",
    Inches(0.7), Inches(1.78), Inches(5.6), Inches(0.5),
    size=12, color=WHITE)

for d, color in [
    ("High cognitive load with multiple competing goals daily", TERRA),
    ("Tech-comfortable — tolerates API-first early product", GOLD),
    ("Strong word-of-mouth: hostels, Discord, WhatsApp groups", DUSK),
    ("Measurable outcomes: completion rate, stress, exam scores", TERRA),
    ("Includes students with ADHD — Jarvis externalises the planning\nbrain, reducing executive function burden significantly", GOLD),
]:
    bullet(s6, color, d,
           Inches(0.7),
           Inches(2.38) + [
               "High cognitive load with multiple competing goals daily",
               "Tech-comfortable — tolerates API-first early product",
               "Strong word-of-mouth: hostels, Discord, WhatsApp groups",
               "Measurable outcomes: completion rate, stress, exam scores",
               "Includes students with ADHD — Jarvis externalises the planning\nbrain, reducing executive function burden significantly",
           ].index(d) * Inches(0.46), size=11)

# Right — why now
card_block(s6, Inches(6.8), Inches(1.38), Inches(6.0), Inches(3.6))
txt(s6, "WHY NOW",
    Inches(7.0), Inches(1.5), Inches(5), Inches(0.25),
    size=8, bold=True, color=TERRA)

for i, (title, body, color) in enumerate([
    ("Local AI is finally real",      "On-device LLMs now run at production quality. Local-first AI wasn't viable 18 months ago.", TERRA),
    ("LLM + solver = trustworthy",    "Combining constraint solvers with language models produces plans users can trust. The stack just matured.", GOLD),
    ("Gen Z burnout is peaking",      "Post-pandemic student burnout at record highs. Demand for low-friction intelligent tools has never been higher.", DUSK),
]):
    t_top = Inches(1.85) + i * Inches(1.1)
    txt(s6, title, Inches(7.0), t_top, Inches(5.5), Inches(0.28),
        size=12, bold=True, color=color)
    txt(s6, body,  Inches(7.0), t_top + Inches(0.3), Inches(5.5), Inches(0.65),
        size=11, color=MUTED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — VISION
# ════════════════════════════════════════════════════════════════════════════
s7 = add_slide(); bg(s7); hud(s7); scan(s7)

tag(s7, "06 / VISION", Inches(0.6), Inches(0.32))
txt(s7, "The operating system for human potential.",
    Inches(0.6), Inches(0.65), Inches(11), Inches(0.6),
    size=30, bold=True, color=WHITE)
txt(s7, "If successful, Jarvis becomes the layer between what you want to achieve and how you spend your time.",
    Inches(0.6), Inches(1.3), Inches(11.5), Inches(0.35),
    size=13, color=MUTED)

phases = [
    ("Phase 1  ·  2026", TERRA,
     "Student Productivity",
     "First 100 → 10K users\nEngineering students, India + SEA\nFreemium, local-first, word-of-mouth"),
    ("Phase 2  ·  2027", GOLD,
     "Professional Expansion",
     "Knowledge workers, freelancers\nTeam scheduling & project layer\nPremium tier, API access"),
    ("Phase 3  ·  2028+", DUSK,
     "Enterprise Platform",
     "Org-wide cognitive load management\nEnterprise contracts + API licensing\n~$90/mo per 1K daily active users"),
]

for i, (phase, color, title, body) in enumerate(phases):
    l = Inches(0.5) + i * Inches(4.22)
    card_block(s7, l, Inches(1.85), Inches(4.02), Inches(2.6), color)
    rect(s7, l, Inches(1.85), Inches(0.06), Inches(2.6), color)
    txt(s7, phase, l+Inches(0.2), Inches(1.98), Inches(3.7), Inches(0.28),
        size=10, bold=True, color=color)
    txt(s7, title, l+Inches(0.2), Inches(2.3),  Inches(3.7), Inches(0.32),
        size=14, bold=True, color=WHITE)
    txt(s7, body,  l+Inches(0.2), Inches(2.68), Inches(3.7), Inches(1.6),
        size=11, color=MUTED)

# TAM/SAM/SOM
rect(s7, Inches(0.5), Inches(4.65), Inches(12.33), Inches(0.95), CARD)
for i, (lbl, val, color) in enumerate([
    ("TAM", "$14.3B  Global productivity software market", MUTED),
    ("SAM", "$2.1B   AI-powered personal productivity", GOLD),
    ("SOM", "$140M  Student segment · India + SEA · 3 year", TERRA),
]):
    col_l = Inches(0.8) + i * Inches(4.2)
    txt(s7, lbl, col_l, Inches(4.78), Inches(0.55), Inches(0.28),
        size=9, bold=True, color=color)
    txt(s7, val, col_l + Inches(0.6), Inches(4.78), Inches(3.5), Inches(0.28),
        size=11, color=WHITE)

txt(s7, "Core moat: every interaction trains PEARL. The behavioral dataset is proprietary. "
        "Switching cost compounds with time — the longer you use Jarvis, the more it knows you.",
    Inches(0.6), Inches(5.82), Inches(11.5), Inches(0.38),
    size=12, bold=True, color=TERRA)


# ════════════════════════════════════════════════════════════════════════════
# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TEAM
# ════════════════════════════════════════════════════════════════════════════
s8 = add_slide(); bg(s8); hud(s8, color=GOLD); scan(s8, color=GOLD)

tag(s8, "THE TEAM", Inches(0.6), Inches(0.32), color=GOLD)
txt(s8, "Built by people who lived the problem.",
    Inches(0.6), Inches(0.65), Inches(11), Inches(0.6),
    size=30, bold=True, color=WHITE)

team = [
    ("Murali Madhav",   TERRA, "Founder & CEO",
     "Engineering student, full-stack & AI developer.\nBuilt Jarvis Engine from scratch — backend, scheduler, memory architecture, LLM routing.\nLives the problem every day."),
    ("Rai Yadav",       GOLD,  "Co-founder",
     "Add role / background here.\nKey contribution to the project.\nWhat they bring to the team."),
    ("Karthik Mehra",   DUSK,  "Co-founder",
     "Add role / background here.\nKey contribution to the project.\nWhat they bring to the team."),
]

card_w = Inches(4.0)
card_h = Inches(3.8)
for i, (name, color, role, bio) in enumerate(team):
    l = Inches(0.5) + i * (card_w + Inches(0.16))
    card_block(s8, l, Inches(1.55), card_w, card_h, color)

    # Avatar circle placeholder
    circ = rect(s8, l + Inches(0.2), Inches(1.75), Inches(0.65), Inches(0.65), color)

    txt(s8, name, l+Inches(1.0), Inches(1.78),
        card_w-Inches(1.1), Inches(0.38), size=15, bold=True, color=WHITE)
    txt(s8, role, l+Inches(1.0), Inches(2.15),
        card_w-Inches(1.1), Inches(0.28), size=10, bold=True, color=color)

    # Divider
    rect(s8, l+Inches(0.2), Inches(2.55), card_w-Inches(0.35), Inches(0.02), DIMMED)

    txt(s8, bio, l+Inches(0.2), Inches(2.7),
        card_w-Inches(0.35), Inches(2.4), size=11, color=MUTED)

txt(s8, "We are not just building a product — we are building the infrastructure for how the next generation manages their minds.",
    Inches(0.6), Inches(5.55), Inches(11.5), Inches(0.38),
    size=12, bold=True, color=WHITE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — INTERVIEW EVIDENCE (Appendix)
# ════════════════════════════════════════════════════════════════════════════
s9 = add_slide(); bg(s9); hud(s9, color=DUSK); scan(s9, color=DUSK)

tag(s9, "APPENDIX / INTERVIEW EVIDENCE", Inches(0.6), Inches(0.32), color=DUSK)
txt(s9, "5 Interviews — Problem Validation",
    Inches(0.6), Inches(0.65), Inches(9), Inches(0.5),
    size=26, bold=True, color=WHITE)
txt(s9, "Conducted March 2026  ·  IIT Madras · IIT KGP · IIT Delhi · Zuvees · Transient AI  ·  Full transcripts attached.",
    Inches(0.6), Inches(1.2), Inches(12), Inches(0.3),
    size=11, color=MUTED, italic=True)

# Table header
cols = [
    ("#",           Inches(0.35), Inches(0.5)),
    ("Profile",     Inches(2.4),  Inches(0.9)),
    ("Key Quote",   Inches(6.5),  Inches(3.35)),
    ("Pain Signal", Inches(2.4),  Inches(9.9)),
]
row_h  = Inches(0.78)
hdr_t  = Inches(1.65)

for lbl, w, l in cols:
    rect(s9, l, hdr_t, w, Inches(0.32), DUSK)
    txt(s9, lbl, l+Inches(0.08), hdr_t+Inches(0.05), w, Inches(0.25),
        size=9, bold=True, color=CANVAS)

rows = [
    ("1", "Nikhilesh\n3rd yr · IIT Madras",
     "I have a to-do list, a calendar, and Notion.\nNone of them tell me what to actually do right now.",
     "Planning paralysis"),
    ("2", "Tharun\n2nd yr · IIT Kharagpur",
     "The guilt of seeing undone tasks is worse than\nnot having a plan. I restart from zero every week.",
     "Guilt / missed tasks"),
    ("3", "Vignesh\n3rd yr · IIT Delhi",
     "Every app solves one problem and creates two more.\nYou end up managing the system, not your work.",
     "Tool overload"),
    ("4", "Abhishek Dhaiya\nCo-founder, Zuvees",
     "The gap isn't features — it's cognitive load absorption.\nThat's what Jarvis is built to do.",
     "Market validation"),
    ("5", "Sreekesh Menon\nCo-founder, Transient AI",
     "Local AI, constraint solvers, peak burnout — all converging.\nAnyone who builds this layer now has a 3-year head start.",
     "Timing / Why Now"),
]

for ri, row_data in enumerate(rows):
    t = Inches(2.0) + ri * row_h
    row_bg = CARD if ri % 2 == 0 else CARD2
    for cell, w, l in zip(row_data, [c[1] for c in cols], [c[2] for c in cols]):
        rect(s9, l, t, w, row_h - Inches(0.03), row_bg)
        txt(s9, cell, l+Inches(0.08), t+Inches(0.08),
            w-Inches(0.12), row_h-Inches(0.1),
            size=10, color=MUTED if ri == 99 else WHITE)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "Jarvis_VC_Pitch_v6.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
